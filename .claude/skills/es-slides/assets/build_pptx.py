# -*- coding: utf-8 -*-
"""
ES（エントリーシート）→ 就活スライド（.pptx）ジェネレーター
=============================================================

3レイアウト（①自己紹介／②志望動機／③自己PR）を 16:9 の .pptx に書き出す。

使い方:
    pip install python-pptx
    python build_pptx.py            # 既定の出力先 ./es-slides.pptx に保存
    python build_pptx.py out.pptx   # 出力先を指定

編集するのは基本的にこのファイル上部の THEME（配色）と DATA（本文）だけ。
本文中の **強調** は、テンプレートのアクセントカラー（青）＋太字になる。
例: "作業時間を約**6割削減**した経験" → 「6割削減」だけ青太字。

デザイン仕様・文字数ルールは references/design-system.md, content-guide.md を参照。
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ============================================================
# 1) THEME — 配色とフォント（テンプレートに合わせて差し替える）
# ============================================================
THEME = {
    "navy":   "#173D63",  # 見出しボックス／タイトル文字
    "navy2":  "#1F5C8B",  # （予備）濃色バリエーション
    "blue":   "#1F7FC0",  # アクセント（数値・キーワード・**強調**）
    "line":   "#2E6FA8",  # タイトル下線
    "ink":    "#2B2F36",  # 本文
    "muted":  "#6B7683",  # 補足ラベル（大学名・英字など）
    "card":   "#FFFFFF",  # カード背景
    "card_bd":"#E3E9F0",  # カード枠線
    "page":   "#EEF3F8",  # ページ背景
    "poly":   "#DCE8F4",  # 隅のローポリ風アクセント
    "font":   "Yu Gothic",
}

# ============================================================
# 2) SIZES — 文字サイズ(pt)。全体を大きく/小さくしたいときはここを触る
# ============================================================
SIZES = {
    "title": 24, "name": 30, "name_en": 13, "aff": 14, "tag": 13,
    "hd": 16, "body": 14, "label": 13, "headline": 26, "note": 15,
    "footer": 10,
}

# ============================================================
# 3) DATA — 本文（ここを自分のESの内容に差し替える）
#    "**...**" で囲むと青太字（アクセント）になる。
# ============================================================
DATA = {
    "intro": {
        "affiliation": "○○大学 ○○学部 ○○学科　4年",
        "name": "田中 太郎",
        "name_en": "Taro Tanaka",
        "tag": "社会学 専攻",
        "cards": [
            {"title": "研究内容", "bullets": [
                "地域コミュニティにおける**情報格差**の問題をテーマに研究",
                "「技術が人にどう使われるか」という視点で調査・分析",
            ]},
            {"title": "研究を通して感じたこと", "bullets": [
                "ITツールが人々の生活を大きく変える可能性を実感",
                "「使う側」から**「作る側」**へ回りたいと考えるように",
            ]},
            {"title": "現在行っていること", "bullets": [
                "独学で**プログラミングの基礎**を学習中",
                "ExcelマクロやPythonで身近な作業の自動化に挑戦",
            ]},
        ],
    },
    "motivation": {
        "label": "POINT ─ ひとことで言うと",
        # 見出し（大きい文字）。行ごとにリストで。** でアクセント。
        "headline": ["“**人に寄り添うものづくり**” に、", "強く共感したからです。"],
        "note_title": "動機の補足",
        "note": [
            "文系で培った人間理解の視点を活かせるから。",
            "非エンジニアの視点を強みに、使いやすいシステムづくりに貢献したいと考えています。",
        ],
    },
    "self_pr": {
        "label": "MY STRENGTH ─ 私の強み",
        "headline": ["課題を分解し、**周囲を巻き込みながら**", "解決に導く力"],
        "cards": [
            {"title": "強みの根拠", "note": [
                "集計自動化で作業時間を約**6割削減**した経験。",
                "工程を分解し優先度順に着手、チームを巻き込んで改善しました。",
            ]},
            {"title": "入社後の活かし方", "note": [
                "複雑な要件を構造化し開発現場で活かす。",
                "対話で周囲を巻き込み、使いやすいシステムづくりに貢献します。",
            ]},
        ],
    },
    "footer": {"intro": "Self Introduction", "motivation": "Motivation", "self_pr": "Self Promotion"},
}

# ============================================================
# ---- 以下、描画エンジン（通常は編集不要） ----
# ============================================================
def C(key):
    h = THEME[key].lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = THEME["font"]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def set_run_font(run, size, color, bold=False):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = FONT
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set("typeface", FONT)


def no_shadow(shape):
    sp = shape._element.spPr
    if sp.find(qn('a:effectLst')) is None:
        sp.append(sp.makeelement(qn('a:effectLst'), {}))


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    no_shadow(sp)
    return sp


def parse_runs(text, size, base, accent=None, bold_base=False):
    """'**強調**' を accent色+太字に分割して runs [(text,size,color,bold)] を返す。"""
    accent = accent if accent is not None else C("blue")
    out, is_accent = [], False
    for part in text.split("**"):
        if part:
            out.append((part, size, accent if is_accent else base, True if is_accent else bold_base))
        is_accent = not is_accent
    return out or [("", size, base, bold_base)]


def textbox(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            space_after=6, line_spacing=1.2, shrink=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    if shrink:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, s, col, b) in runs:
            r = p.add_run(); r.text = t; set_run_font(r, s, col, b)
    return tb


def title_block(slide, label):
    textbox(slide, 0.75, 0.42, 11.8, 0.6,
            [[("▶ ", 15, C("blue"), True), (label, SIZES["title"], C("navy"), True)]],
            anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, 0.75, 1.02, 11.83, 0.028, fill=C("line"))


def footer(slide, en):
    rect(slide, 0.75, 7.02, 0.16, 0.16, fill=C("navy"), shape=MSO_SHAPE.OVAL)
    textbox(slide, 0.98, 6.96, 5.0, 0.3, [[(en, SIZES["footer"], C("muted"), False)]],
            anchor=MSO_ANCHOR.MIDDLE)


def base_slide(title, en):
    s = prs.slides.add_slide(BLANK)
    rect(s, -0.05, -0.05, 13.45, 7.6, fill=C("page"))
    tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                             Inches(9.7), Inches(-0.4), Inches(4.3), Inches(2.6))
    tri.fill.solid(); tri.fill.fore_color.rgb = C("poly")
    tri.line.fill.background(); tri.rotation = 180; no_shadow(tri)
    title_block(s, title); footer(s, en)
    return s


def header_card(slide, x, y, w, hd_text, body_h):
    hd_h = 0.55
    rect(slide, x, y, w, hd_h, fill=C("navy"))
    textbox(slide, x, y, w, hd_h, [[(hd_text, SIZES["hd"], WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(slide, x, y + hd_h, w, body_h, fill=C("card"), line=C("card_bd"), line_w=1.0)
    return hd_h


# ---------- ① 自己紹介 ----------
def build_intro():
    d = DATA["intro"]
    s = base_slide("自己紹介", DATA["footer"]["intro"])
    rect(s, 0.75, 1.28, 11.83, 1.15, fill=C("card"), line=C("card_bd"), line_w=1.0)
    rect(s, 0.75, 1.28, 0.09, 1.15, fill=C("navy"))
    textbox(s, 1.02, 1.40, 8.5, 0.4, [[(d["affiliation"], SIZES["aff"], C("muted"), False)]])
    textbox(s, 1.00, 1.72, 8.5, 0.6,
            [[(d["name"], SIZES["name"], C("navy"), True),
              ("  " + d["name_en"], SIZES["name_en"], C("muted"), False)]])
    rect(s, 10.75, 1.66, 1.75, 0.44, fill=C("navy"), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, 10.75, 1.66, 1.75, 0.44, [[(d["tag"], SIZES["tag"], WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    cy, cbody_h, gap = 2.70, 3.95, 0.22
    cw = (11.83 - 2 * gap) / 3
    for i, card in enumerate(d["cards"]):
        x = 0.75 + i * (cw + gap)
        header_card(s, x, cy, cw, card["title"], cbody_h)
        paras = [[("●  ", 11, C("blue"), True)] + parse_runs(b, SIZES["body"], C("ink"))
                 for b in card["bullets"]]
        textbox(s, x + 0.15, cy + 0.73, cw - 0.30, cbody_h - 0.35, paras,
                space_after=10, line_spacing=1.25, shrink=True)


# ---------- ② 志望動機 ----------
def build_motivation():
    d = DATA["motivation"]
    s = base_slide("志望動機", DATA["footer"]["motivation"])
    ly, lh = 1.30, 2.05
    rect(s, 0.75, ly, 11.83, lh, fill=C("card"), line=C("card_bd"), line_w=1.0)
    rect(s, 0.75, ly, 0.10, lh, fill=C("blue"))
    textbox(s, 1.05, ly + 0.22, 11.3, 0.4, [[(d["label"], SIZES["label"], C("blue"), True)]])
    textbox(s, 1.03, ly + 0.72, 11.4, 1.2,
            [parse_runs(line, SIZES["headline"], C("navy")) for line in d["headline"]],
            line_spacing=1.3, space_after=2)
    ny = ly + lh + 0.28
    header_card(s, 0.75, ny, 11.83, d["note_title"], 6.7 - (ny + 0.55))
    textbox(s, 1.05, ny + 0.72, 11.2, 1.5,
            [parse_runs(line, SIZES["note"], C("ink")) for line in d["note"]],
            line_spacing=1.5, space_after=8, shrink=True)


# ---------- ③ 自己PR ----------
def build_self_pr():
    d = DATA["self_pr"]
    s = base_slide("自己PR", DATA["footer"]["self_pr"])
    ly, lh = 1.30, 1.95
    rect(s, 0.75, ly, 11.83, lh, fill=C("card"), line=C("card_bd"), line_w=1.0)
    rect(s, 0.75, ly, 0.10, lh, fill=C("blue"))
    textbox(s, 1.05, ly + 0.20, 11.3, 0.4, [[(d["label"], SIZES["label"], C("blue"), True)]])
    textbox(s, 1.03, ly + 0.66, 11.4, 1.2,
            [parse_runs(line, SIZES["headline"], C("navy")) for line in d["headline"]],
            line_spacing=1.3, space_after=2)
    sy = ly + lh + 0.28
    sbody_h = 6.7 - (sy + 0.55)
    sgap = 0.28
    sw = (11.83 - sgap) / 2
    for i, card in enumerate(d["cards"]):
        x = 0.75 + i * (sw + sgap)
        header_card(s, x, sy, sw, card["title"], sbody_h)
        paras = [parse_runs(line, SIZES["note"], C("ink")) for line in card["note"]]
        textbox(s, x + 0.28, sy + 0.70, sw - 0.55, sbody_h - 0.30, paras,
                line_spacing=1.4, space_after=8, anchor=MSO_ANCHOR.MIDDLE, shrink=True)


build_intro()
build_motivation()
build_self_pr()

out = sys.argv[1] if len(sys.argv) > 1 else "es-slides.pptx"
prs.save(out)
print("saved", out)
